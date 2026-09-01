"""Generate synthetic test fixtures.

Every fixture is authored here from scratch, so the repository contains no
third-party or copyrighted documents and the expected content of each file is
known exactly - which is what makes it possible to assert on conversion
*quality* rather than merely on a 200 response.

The OCR fixtures render text through HarfBuzz so Bengali is shaped correctly
(see ``scripts/textrender.py``). Rendering Bengali without shaping produces
images whose glyphs are in the wrong order, and OCR run against those would be
measuring the renderer, not the recogniser.

Usage::

    python scripts/generate_fixtures.py [--out fixtures/generated]
"""

from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw, ImageFont

sys.path.insert(0, str(Path(__file__).resolve().parent))
from textrender import draw_text, text_width  # noqa: E402

# --- fonts ------------------------------------------------------------------

LATIN_FONT_CANDIDATES = [
    "C:/Windows/Fonts/arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
]
LATIN_BOLD_CANDIDATES = [
    "C:/Windows/Fonts/arialbd.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
]
BENGALI_FONT_CANDIDATES = [
    "C:/Windows/Fonts/kalpurush.ttf",
    "C:/Windows/Fonts/Nirmala.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansBengali-Regular.ttf",
    "/usr/share/fonts/truetype/lohit-bengali/Lohit-Bengali.ttf",
]


def _first_font(candidates: list[str], label: str) -> str:
    for candidate in candidates:
        if Path(candidate).is_file():
            return candidate
    raise SystemExit(
        f"No {label} font found. Looked for:\n  " + "\n  ".join(candidates)
    )


LATIN = _first_font(LATIN_FONT_CANDIDATES, "Latin")
LATIN_BOLD = _first_font(LATIN_BOLD_CANDIDATES, "Latin bold")
BENGALI = _first_font(BENGALI_FONT_CANDIDATES, "Bengali")

PAGE = (1700, 2200)
WHITE = 255


# --- expected content -------------------------------------------------------

BENGALI_LINES = [
    "বার্ষিক প্রতিবেদন ২০২৬",
    "নির্বাহী সারসংক্ষেপ",
    "এই প্রতিবেদনে বছরের চারটি প্রান্তিকের",
    "আর্থিক ফলাফল উপস্থাপন করা হয়েছে।",
    "শিক্ষা ও স্বাস্থ্য খাতে বিনিয়োগ বেড়েছে।",
]

MIXED_LINES = [
    ("Annual Report 2026", LATIN_BOLD, 56),
    ("বার্ষিক প্রতিবেদন", BENGALI, 48),
    ("Revenue grew by 18 percent this year.", LATIN, 30),
    ("রাজস্ব বেড়েছে আঠারো শতাংশ।", BENGALI, 30),
]


# --- helpers ----------------------------------------------------------------


def _blank() -> Image.Image:
    return Image.new("L", PAGE, WHITE)


def _latin(image: Image.Image, xy, text: str, size: int, bold: bool = False) -> None:
    draw = ImageDraw.Draw(image)
    font = ImageFont.truetype(LATIN_BOLD if bold else LATIN, size)
    draw.text(xy, text, font=font, fill=0)


def _save_png(image: Image.Image, path: Path) -> None:
    image.save(path)


def _save_scanned_pdf(images: list[Image.Image], path: Path) -> None:
    """Write image-only pages: no text layer, so OCR is the only way in."""
    rgb = [img.convert("RGB") for img in images]
    rgb[0].save(
        path,
        save_all=True,
        append_images=rgb[1:],
        # 200 DPI gives a realistic page size rather than a 23-inch monster.
        resolution=200.0,
    )


# --- image / scanned fixtures -----------------------------------------------


def build_english_page() -> Image.Image:
    image = _blank()
    _latin(image, (100, 90), "Annual Report 2026", 64, bold=True)
    _latin(image, (100, 230), "Executive Summary", 44, bold=True)
    for index, line in enumerate(
        [
            "This report summarises the financial performance of the",
            "company across the four quarters of the year, including",
            "segment revenue and operating expenses.",
        ]
    ):
        _latin(image, (100, 320 + index * 44), line, 28)

    _latin(image, (100, 500), "Key Highlights", 44, bold=True)
    for index, line in enumerate(
        [
            "- Revenue grew by 18 percent year over year",
            "- Operating margin improved to 24 percent",
            "- Headcount increased from 400 to 520",
        ]
    ):
        _latin(image, (100, 590 + index * 44), line, 28)

    _latin(image, (100, 770), "Quarterly Revenue", 44, bold=True)
    rows = [
        ("Quarter", "Revenue", "Growth"),
        ("Q1", "50000", "12%"),
        ("Q2", "61000", "18%"),
        ("Q3", "72000", "21%"),
        ("Q4", "88000", "24%"),
    ]
    y = 870
    for row in rows:
        _latin(image, (100, y), row[0], 28)
        _latin(image, (620, y), row[1], 28)
        _latin(image, (1060, y), row[2], 28)
        y += 54
    return image


def build_bengali_page() -> Image.Image:
    image = _blank()
    draw_text(image, (100, 100), BENGALI_LINES[0], BENGALI, 60)
    draw_text(image, (100, 260), BENGALI_LINES[1], BENGALI, 44)
    for index, line in enumerate(BENGALI_LINES[2:]):
        draw_text(image, (100, 380 + index * 70), line, BENGALI, 32)
    return image


def build_mixed_page() -> Image.Image:
    image = _blank()
    y = 110
    for text, font, size in MIXED_LINES:
        if font is BENGALI:
            draw_text(image, (100, y), text, font, size)
        else:
            _latin(image, (100, y), text, size, bold=(font == LATIN_BOLD))
        y += int(size * 2.2)
    return image


def build_table_page() -> Image.Image:
    image = _blank()
    _latin(image, (100, 100), "Inventory Report", 52, bold=True)
    rows = [
        ("Item", "Warehouse", "Quantity", "Value"),
        ("Widget", "North", "1200", "48000"),
        ("Gadget", "South", "340", "17000"),
        ("Sprocket", "East", "875", "26250"),
        ("Flange", "West", "60", "9000"),
    ]
    y = 260
    for row in rows:
        for column, x in zip(row, (100, 520, 950, 1320)):
            _latin(image, (x, y), column, 30)
        y += 58
    return image


def build_multicolumn_page() -> Image.Image:
    """Two columns under a spanning headline, to exercise reading order."""
    image = _blank()
    _latin(image, (100, 90), "Research Digest", 56, bold=True)

    left = [
        "The first column begins here and",
        "continues down the page with",
        "several lines of body text that",
        "should be read before anything",
        "in the second column is reached.",
        "Reading order matters greatly.",
    ]
    right = [
        "The second column starts at the",
        "top of the page again, to the",
        "right of the gutter, and must",
        "follow the first column in the",
        "converted Markdown output.",
        "Order is verified by tests.",
    ]
    for index, line in enumerate(left):
        _latin(image, (100, 260 + index * 52), line, 28)
    for index, line in enumerate(right):
        _latin(image, (950, 260 + index * 52), line, 28)
    return image


def build_skewed_page() -> Image.Image:
    return build_english_page().rotate(
        -6.0, resample=Image.Resampling.BICUBIC, fillcolor=WHITE
    )


def build_low_quality_page() -> Image.Image:
    """Faded, speckled scan - the case preprocessing has to rescue."""
    image = _blank()
    _latin(image, (100, 120), "Faded Receipt", 54, bold=True)
    for index, line in enumerate(
        [
            "This page simulates a poor quality scan with",
            "low contrast and sensor noise throughout.",
            "Total amount due is 1450 dollars.",
        ]
    ):
        _latin(image, (100, 280 + index * 56), line, 32)

    array = np.asarray(image, dtype=np.int16)
    # Compress the tonal range towards mid-grey, then add speckle. Tuned to a
    # genuinely poor scan that a good pipeline can still recover - not to an
    # unreadable one, which would only assert that OCR has limits.
    array = (array * 0.55 + 105).astype(np.int16)
    rng = np.random.default_rng(11)
    array += rng.integers(-14, 14, array.shape)
    mask = rng.random(array.shape) < 0.008
    array[mask] = rng.integers(0, 255, int(mask.sum()))
    return Image.fromarray(np.clip(array, 0, 255).astype(np.uint8), "L")


# --- document fixtures ------------------------------------------------------


def build_text_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    pdf = canvas.Canvas(str(path), pagesize=letter)
    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(72, 720, "Quarterly Update")
    pdf.setFont("Helvetica", 11)
    for index, line in enumerate(
        [
            "This document has a real text layer embedded in the PDF.",
            "It must be extracted directly, without invoking OCR.",
            "Doing so is both faster and more accurate than recognition.",
        ]
    ):
        pdf.drawString(72, 690 - index * 16, line)
    pdf.showPage()

    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(72, 720, "Page Two")
    pdf.setFont("Helvetica", 11)
    pdf.drawString(72, 695, "Content continues onto a second page.")
    pdf.showPage()
    pdf.save()


def build_docx(path: Path) -> None:
    import docx

    document = docx.Document()
    document.add_heading("Project Proposal", level=1)
    document.add_paragraph("An introductory paragraph describing the project.")
    document.add_heading("Objectives", level=2)
    for item in ("Reduce processing time", "Improve accuracy", "Lower cost"):
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Budget", level=2)
    table = document.add_table(rows=3, cols=2)
    table.style = "Table Grid"
    data = [("Category", "Amount"), ("Engineering", "120000"), ("Research", "45000")]
    for row_index, row in enumerate(data):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = value

    paragraph = document.add_paragraph()
    paragraph.add_run("Bold text").bold = True
    paragraph.add_run(" and ")
    paragraph.add_run("italic text").italic = True
    document.save(path)


def build_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Market Overview"
    title_slide.placeholders[1].text = "Annual strategy review"

    bullet_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    bullet_slide.shapes.title.text = "Key Findings"
    frame = bullet_slide.placeholders[1].text_frame
    frame.text = "Demand increased in every region"
    frame.add_paragraph().text = "Supply constraints eased in the third quarter"
    frame.add_paragraph().text = "Pricing remained stable"
    bullet_slide.notes_slide.notes_text_frame.text = (
        "Remember to mention the supply chain recovery."
    )
    presentation.save(path)


def build_xlsx(path: Path) -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    sales = workbook.active
    sales.title = "Sales"
    for row in [
        ["Product", "Quantity", "Revenue"],
        ["Product A", 100, 50000],
        ["Product B", 250, 90000],
    ]:
        sales.append(row)

    expenses = workbook.create_sheet("Expenses")
    for row in [["Category", "Amount"], ["Marketing", 12000], ["Travel", 4300]]:
        expenses.append(row)
    workbook.save(path)


def build_epub(path: Path) -> None:
    """A minimal but valid EPUB 3, assembled by hand to avoid a dependency."""
    container = """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles><rootfile full-path="OEBPS/content.opf"
   media-type="application/oebps-package+xml"/></rootfiles>
</container>"""

    opf = """<?xml version="1.0" encoding="utf-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="id">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:identifier id="id">urn:uuid:markitdown-fixture</dc:identifier>
    <dc:title>A Short Book</dc:title>
    <dc:language>en</dc:language>
  </metadata>
  <manifest>
    <item id="c1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
    <item id="c2" href="chapter2.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine><itemref idref="c1"/><itemref idref="c2"/></spine>
</package>"""

    def chapter(number: int, title: str, body: str) -> str:
        return f"""<?xml version="1.0" encoding="utf-8"?>
<html xmlns="http://www.w3.org/1999/xhtml"><head><title>{title}</title></head>
<body><h1>Chapter {number}: {title}</h1><p>{body}</p></body></html>"""

    with zipfile.ZipFile(path, "w") as archive:
        # The mimetype entry must be first and stored uncompressed.
        archive.writestr(
            zipfile.ZipInfo("mimetype"),
            "application/epub+zip",
            compress_type=zipfile.ZIP_STORED,
        )
        archive.writestr("META-INF/container.xml", container)
        archive.writestr("OEBPS/content.opf", opf)
        archive.writestr(
            "OEBPS/chapter1.xhtml",
            chapter(1, "The Beginning", "The opening chapter introduces the topic."),
        )
        archive.writestr(
            "OEBPS/chapter2.xhtml",
            chapter(2, "The Middle", "The second chapter develops the argument."),
        )


def build_data_fixtures(out: Path) -> None:
    (out / "table.csv").write_text(
        "Product,Quantity,Revenue\nAlpha,100,50000\nBeta,250,90000\nGamma,75,31000\n",
        encoding="utf-8",
    )

    (out / "records.json").write_text(
        json.dumps(
            {
                "title": "Quarterly Report",
                "team": {"lead": "Ada", "size": 4},
                "items": [
                    {"name": "alpha", "value": 1, "active": True},
                    {"name": "beta", "value": 2, "active": False},
                ],
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    (out / "catalog.xml").write_text(
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<catalog vendor="Acme">\n'
        '  <book id="1"><title>Deep Work</title><price>29</price></book>\n'
        '  <book id="2"><title>Wide Nets</title><price>35</price></book>\n'
        "</catalog>\n",
        encoding="utf-8",
    )

    (out / "page.html").write_text(
        "<!doctype html><html><head><title>Sample Page</title></head><body>"
        "<h1>Sample Page</h1><p>A paragraph with <strong>bold</strong> text and a "
        '<a href="https://example.com">link</a>.</p>'
        "<ul><li>First item</li><li>Second item</li></ul>"
        "<table><tr><th>Name</th><th>Value</th></tr>"
        "<tr><td>alpha</td><td>1</td></tr></table>"
        "</body></html>\n",
        encoding="utf-8",
    )

    (out / "notebook.ipynb").write_text(
        json.dumps(
            {
                "nbformat": 4,
                "nbformat_minor": 5,
                "metadata": {"kernelspec": {"language": "python"}},
                "cells": [
                    {
                        "cell_type": "markdown",
                        "metadata": {},
                        "source": ["# Analysis\n", "\n", "Introductory text."],
                    },
                    {
                        "cell_type": "code",
                        "execution_count": 1,
                        "metadata": {},
                        "outputs": [
                            {
                                "output_type": "stream",
                                "name": "stdout",
                                "text": ["42\n"],
                            }
                        ],
                        "source": ["print(6 * 7)"],
                    },
                ],
            },
            indent=1,
        ),
        encoding="utf-8",
    )


def build_archive(out: Path) -> None:
    with zipfile.ZipFile(out / "bundle.zip", "w", zipfile.ZIP_DEFLATED) as archive:
        archive.write(out / "table.csv", "data/table.csv")
        archive.write(out / "records.json", "data/records.json")
        archive.writestr("notes/readme.txt", "Unsupported entries are skipped.")


# --- entry point ------------------------------------------------------------


def generate(out: Path) -> list[Path]:
    out.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []

    def record(path: Path) -> None:
        written.append(path)

    # Images (all three raster formats we accept).
    english = build_english_page()
    _save_png(english, out / "ocr-english.png")
    record(out / "ocr-english.png")

    english.convert("RGB").save(out / "ocr-english.jpg", quality=92)
    record(out / "ocr-english.jpg")
    english.convert("RGB").save(out / "ocr-english.webp", quality=92)
    record(out / "ocr-english.webp")

    for name, builder in [
        ("ocr-bengali.png", build_bengali_page),
        ("ocr-mixed.png", build_mixed_page),
        ("ocr-table.png", build_table_page),
        ("ocr-multicolumn.png", build_multicolumn_page),
        ("ocr-skewed.png", build_skewed_page),
        ("ocr-low-quality.png", build_low_quality_page),
    ]:
        _save_png(builder(), out / name)
        record(out / name)

    # Scanned (image-only) PDFs.
    _save_scanned_pdf([build_english_page()], out / "scanned-english.pdf")
    record(out / "scanned-english.pdf")
    _save_scanned_pdf([build_bengali_page()], out / "scanned-bengali.pdf")
    record(out / "scanned-bengali.pdf")
    _save_scanned_pdf(
        [build_english_page(), build_table_page()], out / "scanned-multipage.pdf"
    )
    record(out / "scanned-multipage.pdf")

    # Real documents.
    build_text_pdf(out / "text.pdf")
    record(out / "text.pdf")
    build_docx(out / "proposal.docx")
    record(out / "proposal.docx")
    build_pptx(out / "deck.pptx")
    record(out / "deck.pptx")
    build_xlsx(out / "workbook.xlsx")
    record(out / "workbook.xlsx")
    build_epub(out / "book.epub")
    record(out / "book.epub")

    build_data_fixtures(out)
    for name in ("table.csv", "records.json", "catalog.xml", "page.html", "notebook.ipynb"):
        record(out / name)

    build_archive(out)
    record(out / "bundle.zip")

    return written


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--out",
        type=Path,
        default=Path(__file__).resolve().parents[1] / "fixtures" / "generated",
    )
    args = parser.parse_args()

    written = generate(args.out)
    print(f"Wrote {len(written)} fixtures to {args.out}")
    for path in written:
        print(f"  {path.name:26} {path.stat().st_size:>9,} bytes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
